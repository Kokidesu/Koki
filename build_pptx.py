#!/usr/bin/env python3
"""Build the Epstein / transnational crime presentation deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Palette ----
NAVY   = RGBColor(0x0F, 0x2A, 0x43)   # deep navy
SLATE  = RGBColor(0x33, 0x44, 0x55)
ACCENT = RGBColor(0xC0, 0x8A, 0x2B)   # muted gold
LIGHT  = RGBColor(0xF2, 0xF4, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x6B, 0x76, 0x82)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, text, size=18, color=SLATE, bold=False,
        align=PP_ALIGN.LEFT, font="Times New Roman", anchor=MSO_ANCHOR.TOP, italic=False,
        line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return tb


def bullets(slide, x, y, w, h, items, size=18, color=SLATE, gap=6, bold_lead=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        # support "lead: rest"
        if bold_lead and ": " in item:
            lead, rest = item.split(": ", 1)
            r1 = p.add_run(); r1.text = "•  " + lead + ": "
            r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = NAVY; r1.font.name = "Times New Roman"
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = "Times New Roman"
        else:
            r = p.add_run(); r.text = "•  " + item
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Times New Roman"
    return tb


def content_header(slide, kicker, title):
    bg(slide, WHITE)
    rect(slide, 0, 0, SW, Inches(1.25), NAVY)
    rect(slide, 0, Inches(1.25), SW, Inches(0.06), ACCENT)
    txt(slide, Inches(0.6), Inches(0.18), Inches(11), Inches(0.4),
        kicker, size=13, color=ACCENT, bold=True, font="Times New Roman")
    txt(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.7),
        title, size=26, color=WHITE, bold=True, font="Times New Roman")


# ============ SLIDE 1 — TITLE ============
s = add_slide()
bg(s, NAVY)
rect(s, 0, Inches(3.05), SW, Inches(0.06), ACCENT)
txt(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(0.5),
    "TRANSNATIONAL CRIME & THE LIMITS OF INTERNATIONAL JUSTICE",
    size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.1),
    "What the Epstein Case Reveals",
    size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1.0), Inches(3.35), Inches(11.3), Inches(0.6),
    "A case study in cross-border exploitation, jurisdiction, and accountability",
    size=18, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)
txt(s, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.5),
    "Group Presentation  •  Faculty of Law  •  2026",
    size=13, color=GREY, align=PP_ALIGN.CENTER)

# ============ SLIDE 2 — CENTRAL QUESTION / AGENDA ============
s = add_slide()
content_header(s, "OVERVIEW", "Our Central Question & Roadmap")
rect(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.25), LIGHT)
txt(s, Inches(0.9), Inches(1.78), Inches(11.5), Inches(1.0),
    "Why is transnational sexual exploitation so hard to prosecute — and what do\n"
    "international law and cross-border cooperation need in order to work?",
    size=19, color=NAVY, bold=True, italic=True, line_spacing=1.1)
bullets(s, Inches(0.7), Inches(3.2), Inches(12), Inches(3.6), [
    "Part 1 — Framing & the international scope of the case",
    "Part 2 — The international legal framework (Palermo Protocol, UNTOC, jurisdiction & extradition)",
    "Part 3 — Wealth, mobility, and the problem of \"elite impunity\"",
    "Part 4 — Accountability, transparency & lessons (incl. a Japan comparison)",
], size=19, gap=14, color=SLATE)
txt(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
    "Note: This presentation centres victims and legal systems — not sensational detail.",
    size=13, color=GREY, italic=True)

# ============ SLIDE 3 — PART 1 ============
s = add_slide()
content_header(s, "PART 1  •  FRAMING", "Why This Is an International Story")
bullets(s, Inches(0.7), Inches(1.6), Inches(7.4), Inches(5.4), [
    "Who: Jeffrey Epstein, a financier whose network spanned multiple countries.",
    "2008: A lenient Florida plea deal — state charges, limited accountability.",
    "2019: Federal sex-trafficking charges in New York; Epstein died in custody.",
    "Cross-border recruitment of victims, including minors and foreign nationals.",
    "Offshore wealth and properties across the US, US Virgin Islands and Paris.",
    "A globe-spanning social network of internationally prominent figures.",
], size=18, gap=14, bold_lead=True)
rect(s, Inches(8.4), Inches(1.7), Inches(4.3), Inches(4.9), LIGHT)
txt(s, Inches(8.7), Inches(1.95), Inches(3.8), Inches(0.5),
    "WHY IT'S \"INTERNATIONAL\"", size=14, color=ACCENT, bold=True)
bullets(s, Inches(8.7), Inches(2.5), Inches(3.8), Inches(4.0), [
    "Victims moved across borders",
    "Money held offshore",
    "Multiple jurisdictions involved",
    "Foreign & quasi-diplomatic figures",
    "→ A test for international justice",
], size=16, gap=12, color=SLATE)

# ============ SLIDE 4 — PART 2 ============
s = add_slide()
content_header(s, "PART 2  •  LEGAL FRAMEWORK", "Trafficking as a Transnational Crime")
bullets(s, Inches(0.7), Inches(1.6), Inches(12), Inches(5.4), [
    "UN Palermo Protocol (2000): Defines trafficking in persons and obliges states to criminalise it.",
    "UN Convention against Transnational Organized Crime (UNTOC): The parent framework for cooperation.",
    "Jurisdiction: Which state may prosecute when conduct spans several countries?",
    "Extradition & MLATs: Mutual Legal Assistance Treaties enable evidence-sharing across borders.",
    "Foreign nationals: Ghislaine Maxwell (UK) — convicted 2021 of sex-trafficking conspiracy.",
    "Quasi-diplomatic complications: Prince Andrew settled a US civil case (2022) and later lost titles.",
], size=18, gap=14, bold_lead=True)
txt(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
    "Takeaway: The law exists — the difficulty is coordinating sovereign systems.",
    size=14, color=NAVY, bold=True, italic=True)

# ============ SLIDE 5 — PART 3 ============
s = add_slide()
content_header(s, "PART 3  •  WEALTH & MOBILITY", "How Wealth Enabled \"Elite Impunity\"")
bullets(s, Inches(0.7), Inches(1.6), Inches(7.4), Inches(5.4), [
    "Offshore structures and private wealth obscured assets and ownership.",
    "Property in several jurisdictions enabled mobility and evasion.",
    "Identity documents: an expired passport bore Epstein's photo, a false name, and listed Saudi Arabia as residence.",
    "Contrast: the lenient 2008 outcome vs. the 2019 federal prosecution.",
    "Maxwell's 2021 conviction showed accountability is possible — but slow.",
], size=18, gap=13, bold_lead=True)
rect(s, Inches(8.4), Inches(1.7), Inches(4.3), Inches(4.9), LIGHT)
txt(s, Inches(8.7), Inches(1.95), Inches(3.8), Inches(0.6),
    "COMPARATIVE QUESTION", size=14, color=ACCENT, bold=True)
txt(s, Inches(8.7), Inches(2.6), Inches(3.8), Inches(3.8),
    "Would other legal systems — the UK, France, or Japan — have acted faster, "
    "or slower, against a defendant of comparable wealth and mobility?",
    size=17, color=SLATE, italic=True, line_spacing=1.15)

# ============ SLIDE 6 — PART 4 ============
s = add_slide()
content_header(s, "PART 4  •  ACCOUNTABILITY", "Transparency, Survivors & Lessons")
bullets(s, Inches(0.7), Inches(1.6), Inches(12), Inches(4.6), [
    "2025 document-release wave under US transparency legislation brought new scrutiny.",
    "Survivors' concern: disclosures made it easy to identify victims — but not the enablers.",
    "Lesson: transparency must be designed to protect victims while exposing accountability.",
    "Reform: stronger MLATs, faster extradition, victim-centred disclosure, asset tracing.",
    "Japan comparison: anti-trafficking framework & US State Dept. TIP assessments offer a benchmark.",
], size=18, gap=15, bold_lead=True)
rect(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.75), NAVY)
txt(s, Inches(0.9), Inches(6.48), Inches(11.5), Inches(0.5),
    "Closing: Justice across borders depends less on new laws than on the will to cooperate.",
    size=16, color=WHITE, bold=True, italic=True)

# ============ SLIDE 7 — KEY TAKEAWAYS ============
s = add_slide()
content_header(s, "CONCLUSION", "Four Takeaways")
items = [
    ("01", "Crime crosses borders — justice often doesn't.", "Networks are global; enforcement is national."),
    ("02", "The legal tools exist.", "Palermo & UNTOC provide the framework; coordination is the gap."),
    ("03", "Wealth buys time, not immunity.", "Mobility delayed — but did not prevent — accountability."),
    ("04", "Transparency must protect victims.", "Disclosure should expose enablers, not endanger survivors."),
]
y = 1.65
for num, head, sub in items:
    rect(s, Inches(0.7), Inches(y), Inches(0.85), Inches(1.05), ACCENT)
    txt(s, Inches(0.7), Inches(y+0.18), Inches(0.85), Inches(0.7),
        num, size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(1.75), Inches(y+0.05), Inches(10.8), Inches(0.5),
        head, size=19, color=NAVY, bold=True)
    txt(s, Inches(1.75), Inches(y+0.52), Inches(10.8), Inches(0.5),
        sub, size=15, color=SLATE, italic=True)
    y += 1.25

# ============ SLIDE 8 — REFERENCES / THANK YOU ============
s = add_slide()
bg(s, NAVY)
rect(s, 0, Inches(2.0), SW, Inches(0.06), ACCENT)
txt(s, Inches(1.0), Inches(0.9), Inches(11.3), Inches(0.9),
    "Thank You", size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.5),
    "Questions & Discussion", size=18, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)
txt(s, Inches(1.0), Inches(3.3), Inches(11.3), Inches(0.4),
    "SELECTED SOURCES & INSTRUMENTS", size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
bullets(s, Inches(3.3), Inches(3.8), Inches(7), Inches(2.6), [
    "UN Protocol to Prevent, Suppress and Punish Trafficking in Persons (Palermo Protocol, 2000)",
    "UN Convention against Transnational Organized Crime (UNTOC)",
    "US Dept. of Justice filings: US v. Epstein (2019); US v. Maxwell (2021)",
    "US State Department Trafficking in Persons (TIP) Reports",
    "Contemporary news reporting on 2025 document disclosures",
], size=14, gap=8, color=LIGHT)

prs.save("/home/user/Koki/Epstein_Transnational_Justice.pptx")
print("Saved deck with", len(prs.slides._sldIdLst), "slides")
