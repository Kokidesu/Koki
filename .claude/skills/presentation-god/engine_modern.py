#!/usr/bin/env python3
"""Presentation God — Modern theme (MusePass-style).

A clean, light, Japanese-startup / SaaS pitch aesthetic learned from a real deck:
white canvas, gothic sans Japanese, a single vivid blue accent, generous
whitespace, card rows with coloured summary strips, blue process boxes with
arrows, big "jump-ratio" headings, and a quiet "NEXT ▶ … / page" footer.

Design principles baked in (from JP slide-design best practice):
  1 slide = 1 message · plenty of whitespace · 2–3 colours (dark grey + blue) ·
  high jump ratio (big bold heading, small calm body) · minimal decoration.

Fonts:  Japanese → Noto Sans CJK JP (gothic, the clean modern choice) and
        游ゴシック / "Yu Gothic" in the .pptx.   English → Times New Roman.
"""
from __future__ import annotations

import math
import os
import re
from io import BytesIO

from PIL import Image, ImageDraw, ImageFont

SW_IN, SH_IN = 13.333, 7.5
MARGIN = 0.92

# ---------- palette ----------
WHITE = (0xFF, 0xFF, 0xFF)
INK = (0x1B, 0x20, 0x28)        # headings (near-black)
BODY = (0x44, 0x4D, 0x5A)       # body text (dark grey)
MUTE = (0x8A, 0x93, 0xA1)       # captions / footer
BLUE = (0x2F, 0x6B, 0xF0)       # default primary accent
SOFT = (0xF3, 0xF5, 0xF9)       # card grey (theme-neutral)
HAIR = (0xE2, 0xE6, 0xEC)       # hairlines / card border


def _mix(c, other, t):
    return tuple(int(round(c[i] + (other[i] - c[i]) * t)) for i in range(3))


def _shades(accent):
    """Derive a coherent set of tones from a single accent colour."""
    return {
        "accent": accent,
        "dark": _mix(accent, (0, 0, 0), 0.28),       # title gradient end / strong text
        "box": _mix(accent, (0x55, 0x60, 0x80), 0.30),  # process box fill (muted)
        "band": _mix(accent, WHITE, 0.88),           # light summary strip
    }


# Theme presets — switch the whole deck's colour with one name.
THEMES = {
    "blue": _shades((0x2F, 0x6B, 0xF0)),
    "indigo": _shades((0x49, 0x46, 0xE5)),
    "teal": _shades((0x0E, 0x9F, 0x9A)),
    "green": _shades((0x1F, 0x9D, 0x55)),
    "purple": _shades((0x7C, 0x3A, 0xED)),
    "orange": _shades((0xE8, 0x6A, 0x17)),
    "rose": _shades((0xE1, 0x3D, 0x67)),
    "slate": _shades((0x3D, 0x4E, 0x6B)),
}

# ---------- fonts ----------
_NOTO = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
_NOTOB = "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc"
_JP_REG = (_NOTO, 0) if os.path.exists(_NOTO) else ("/usr/share/fonts/opentype/ipafont-gothic/ipagp.ttf", 0)
_JP_BOLD = (_NOTOB, 0) if os.path.exists(_NOTOB) else _JP_REG
_LAT = {
    "r": "/usr/share/fonts/truetype/liberation/LiberationSerif-Regular.ttf",
    "b": "/usr/share/fonts/truetype/liberation/LiberationSerif-Bold.ttf",
    "i": "/usr/share/fonts/truetype/liberation/LiberationSerif-Italic.ttf",
}
PPTX_JP_FONT = "Yu Gothic"
PPTX_LAT_FONT = "Times New Roman"
_CJK_RE = re.compile(r"[぀-ヿ一-鿿０-ｚ]")
_fc: dict = {}


def _has_cjk(t: str) -> bool:
    return bool(_CJK_RE.search(t))


_MEASURE = ImageDraw.Draw(Image.new("RGB", (10, 10)))


def _fit_size(text: str, max_w_in: float, start: int, min_size: int = 18) -> int:
    """Largest point size (≤ start) at which `text` fits `max_w_in` on one line.

    Width is measured at 72 px/in so it is resolution-independent.
    """
    cjk = _has_cjk(text)
    target_px = max_w_in * 72.0
    size = start
    while size > min_size:
        f = _font(size, True, False, cjk)
        if _MEASURE.textlength(text, font=f) <= target_px:
            break
        size -= 2
    return size


def _font(size: int, bold: bool, italic: bool, cjk: bool):
    if cjk:
        path, idx = (_JP_BOLD if bold else _JP_REG)
    else:
        path, idx = (_LAT["b"] if bold else _LAT["i"] if italic else _LAT["r"]), 0
    k = (path, idx, size)
    if k not in _fc:
        _fc[k] = ImageFont.truetype(path, size, index=idx)
    return _fc[k]


class Slide:
    """A collected list of draw ops; rendered by both PIL and python-pptx."""

    def __init__(self, bg=WHITE):
        self.bg = bg
        self.grad = None        # (top, bottom) for a diagonal gradient bg
        self.el: list = []

    def rect(self, x, y, w, h, c, r=0.0):
        self.el.append(("rect", x, y, w, h, c, r))

    def line(self, x1, y1, x2, y2, c, w):
        self.el.append(("line", x1, y1, x2, y2, c, w))

    def arrow(self, x1, y1, x2, y2, c, w):
        self.el.append(("arrow", x1, y1, x2, y2, c, w))

    def text(self, x, y, w, t, size, c=BODY, bold=False, italic=False, align="l", ls=1.18, h=0.6):
        self.el.append(("text", x, y, w, h, t, size, c, bold, italic, align, ls))

    def bullets(self, x, y, w, items, size, c=BODY, gap=9, h=4.0):
        self.el.append(("bullets", x, y, w, h, items, size, c, gap))

    def image(self, x, y, w, h, img):
        self.el.append(("image", x, y, w, h, img))


class ModernDeck:
    def __init__(self, theme="blue"):
        sh = THEMES.get(theme, THEMES["blue"])
        self.accent = sh["accent"]
        self.dark = sh["dark"]
        self.box = sh["box"]
        self.band = sh["band"]
        self.slides: list[Slide] = []
        self.total = 0

    def _add(self, s):
        self.slides.append(s)
        return s

    # ---------------- slide types ----------------
    def title(self, kicker, headline, subline="", company="", note="", mock=None):
        s = Slide()
        s.grad = (self.accent, self.dark)
        # subtle lighter wedge top-right for depth
        s.text(MARGIN, 1.45, 8.5, kicker, 19, WHITE, bold=True)
        # headline: large, bold, left-aligned, multi-line allowed via \n
        s.text(MARGIN, 2.25, 8.6, headline, 48, WHITE, bold=True, ls=1.16, h=3.2)
        if subline:
            s.text(MARGIN, 5.6, 8.4, subline, 21, (0xDF, 0xE8, 0xFF), ls=1.3, h=0.9)
        if company:
            s.text(MARGIN, 6.75, 6.0, company, 16, (0xCF, 0xDD, 0xFF), bold=True)
        if note:
            s.text(7.0, 6.98, 5.4, note, 13, (0xBF, 0xD0, 0xFB), align="r")
        if mock is not None:
            s.image(9.2, 1.2, 3.2, 5.1, mock)
        return self._add(s)

    def _head(self, s, section, title, lead=""):
        # thin accent bar + small section label
        s.rect(MARGIN, 0.62, 0.07, 0.42, self.accent)
        s.text(MARGIN + 0.2, 0.60, 10.5, section, 15, MUTE, bold=True, h=0.4)
        s.text(MARGIN, 1.18, 11.5, title, 35, INK, bold=True, h=1.0, ls=1.08)
        if lead:
            s.text(MARGIN, 2.22, 11.5, lead, 18, BODY, ls=1.3, h=0.8)

    def section(self, no, total, title, lead=""):
        s = Slide()
        s.grad = (self.accent, self.dark)
        s.text(MARGIN, 2.5, 4.0, f"{no:02d} / {total:02d}", 18, (0xCF, 0xDD, 0xFF), bold=True)
        s.rect(MARGIN, 3.15, 1.7, 0.06, WHITE)
        s.text(MARGIN, 3.42, 10.8, title, 40, WHITE, bold=True, h=1.7, ls=1.1)
        if lead:
            s.text(MARGIN, 5.25, 10.0, lead, 20, (0xDF, 0xE8, 0xFF), ls=1.3, h=1.0)
        return self._add(s)

    def cards(self, section, title, cards, lead="", next_=""):
        """cards = list of (heading, [bullets], summary). 2–4 cards."""
        s = Slide()
        self._head(s, section, title, lead)
        n = len(cards)
        gap = 0.4
        total_w = SW_IN - 2 * MARGIN
        cw = (total_w - gap * (n - 1)) / n
        top = 3.05
        ch = 3.55
        for i, (head, bul, summ) in enumerate(cards):
            x = MARGIN + i * (cw + gap)
            s.rect(x, top, cw, ch, SOFT, r=0.12)
            s.rect(x, top, cw, 0.09, self.accent, r=0.0)  # top accent line
            s.text(x + 0.3, top + 0.34, cw - 0.6, head, 22, INK, bold=True, h=0.8, ls=1.1)
            s.bullets(x + 0.3, top + 1.32, cw - 0.6, bul, 16, BODY, gap=10)
            # summary strip at the bottom of the card
            strip_h = 0.92
            s.rect(x, top + ch - strip_h, cw, strip_h, self.band, r=0.0)
            s.text(x + 0.3, top + ch - strip_h + 0.2, cw - 0.6, summ, 16, self.dark, bold=True, h=0.6, ls=1.15)
        self._foot(s, next_)
        return self._add(s)

    def points(self, section, title, items, lead="", next_=""):
        """A single clean list — 1 message, big whitespace."""
        s = Slide()
        self._head(s, section, title, lead)
        n = len(items)
        y0 = 3.1 if n <= 3 else 2.85
        gap = {1: 34, 2: 28, 3: 22, 4: 15}.get(n, 11)
        s.bullets(MARGIN + 0.1, y0, SW_IN - 2 * MARGIN - 0.2, items, 21, BODY, gap=gap)
        self._foot(s, next_)
        return self._add(s)

    def kpis(self, section, title, stats, lead="", next_=""):
        """stats = list of (number, label). Big jump-ratio numbers in accent."""
        s = Slide()
        self._head(s, section, title, lead)
        n = len(stats)
        gap = 0.4
        total_w = SW_IN - 2 * MARGIN
        cw = (total_w - gap * (n - 1)) / n
        top = 3.05
        ch = 3.0
        for i, (num, lab) in enumerate(stats):
            x = MARGIN + i * (cw + gap)
            s.rect(x, top, cw, ch, SOFT, r=0.12)
            num_size = _fit_size(num, cw - 0.5, start=60, min_size=30)
            s.text(x + 0.2, top + 0.55, cw - 0.4, num, num_size, self.accent, bold=True, align="c", h=1.3)
            s.text(x + 0.2, top + 2.0, cw - 0.4, lab, 16, BODY, align="c", ls=1.25, h=0.85)
        self._foot(s, next_)
        return self._add(s)

    def flow(self, section, title, steps, lead="", next_=""):
        """steps = list of (name, sub). Blue rounded boxes + arrows, 2x2 or row."""
        s = Slide()
        self._head(s, section, title, lead)
        n = len(steps)
        per_row = 2 if n == 4 else min(n, 4)
        gap_x, gap_y = 1.1, 0.7
        rows = math.ceil(n / per_row)
        total_w = SW_IN - 2 * MARGIN
        bw = (total_w - gap_x * (per_row - 1)) / per_row
        bh = 1.7 if rows > 1 else 2.4
        top = 3.1 if rows > 1 else 3.2
        centers = []
        for i, (name, sub) in enumerate(steps):
            r, c = divmod(i, per_row)
            x = MARGIN + c * (bw + gap_x)
            y = top + r * (bh + gap_y)
            s.rect(x, y, bw, bh, self.box, r=0.1)
            s.text(x + 0.25, y + bh / 2 - 0.62, bw - 0.5, name, 23, WHITE, bold=True, align="c", h=0.6)
            s.text(x + 0.25, y + bh / 2 + 0.08, bw - 0.5, sub, 15, (0xE4, 0xEC, 0xFF), align="c", ls=1.2, h=0.7)
            centers.append((x, y, bw, bh))
        # arrows between consecutive boxes
        for i in range(n - 1):
            x0, y0, w0, h0 = centers[i]
            x1, y1, w1, h1 = centers[i + 1]
            if abs(y0 - y1) < 0.01:  # same row → horizontal
                s.arrow(x0 + w0 + 0.12, y0 + h0 / 2, x1 - 0.12, y1 + h1 / 2, BODY, 3)
            else:  # wrap to next row → diagonal
                s.arrow(x0 + w0 / 2, y0 + h0 + 0.1, x1 + w1 / 2, y1 - 0.1, (0xC2, 0xC8, 0xD2), 3)
        self._foot(s, next_)
        return self._add(s)

    def takeaways(self, section, title, items, sources="", next_=""):
        """items = list of (n, head, sub). Numbered, clean."""
        s = Slide()
        self._head(s, section, title)
        y = 2.95
        for num, head, sub in items:
            s.rect(MARGIN, y, 0.7, 0.7, self.accent, r=0.1)
            s.text(MARGIN, y + 0.15, 0.7, num, 23, WHITE, bold=True, align="c", h=0.4)
            s.text(MARGIN + 1.05, y + 0.04, 10.5, head, 21, INK, bold=True, h=0.5)
            s.text(MARGIN + 1.05, y + 0.56, 10.5, sub, 15, BODY, h=0.45, ls=1.2)
            y += 1.12
        if sources:
            s.rect(MARGIN, 6.42, SW_IN - 2 * MARGIN, 0.66, SOFT, r=0.08)
            s.text(MARGIN + 0.25, 6.5, 11.0, "SOURCES / 出典", 11, MUTE, bold=True, h=0.3)
            s.text(MARGIN + 0.25, 6.76, 11.4, sources, 11, BODY, ls=1.2, h=0.3)
        return self._add(s)

    def _foot(self, s, next_):
        if next_:
            s.text(MARGIN, 7.02, 8.0, f"NEXT ▶ {next_}", 12.5, MUTE, h=0.3)

    # ---------------- build ----------------
    def build(self, outdir, name, pdf=True, contact=True):
        os.makedirs(outdir, exist_ok=True)
        self.total = len(self.slides)
        out = {"n": self.total}
        out["pptx"] = os.path.join(outdir, f"{name}.pptx")
        self._pptx(out["pptx"])
        previews = []
        for i, s in enumerate(self.slides):
            p = os.path.join(outdir, f"{name}_{i + 1:02d}.png")
            self._png(s, i + 1, p)
            previews.append(p)
        if contact:
            out["contact"] = os.path.join(outdir, f"{name}_contact.png")
            self._contact(previews, out["contact"])
        if pdf:
            out["pdf"] = os.path.join(outdir, f"{name}.pdf")
            self._pdf(name, outdir, out["pdf"])
        return out

    # ---- PNG ----
    _PX = 110

    def _I(self, v):
        return int(round(v * self._PX))

    def _pt(self, p):
        return int(round(p * self._PX / 72.0))

    def _wrap(self, d, t, f, maxw):
        cjk = _has_cjk(t) and " " not in t.strip()
        units = list(t) if cjk else t.split(" ")
        join = "" if cjk else " "
        out, cur = [], ""
        for w in units:
            cand = (cur + join + w) if cur else w
            if d.textlength(cand, font=f) <= maxw:
                cur = cand
            else:
                if cur:
                    out.append(cur)
                cur = w
        if cur:
            out.append(cur)
        return out or [""]

    def _png(self, s: Slide, pageno: int, path: str):
        W, H = self._I(SW_IN), self._I(SH_IN)
        if s.grad:
            img = self._gradient(W, H, s.grad[0], s.grad[1])
        else:
            img = Image.new("RGB", (W, H), s.bg)
        d = ImageDraw.Draw(img)
        for el in s.el:
            self._draw(d, img, el)
        # page number (not on title = slide 1)
        if pageno > 1 and not s.grad:
            f = _font(self._pt(11), False, False, False)
            d.text((W - self._I(0.7), H - self._I(0.5)), f"{pageno:02d}", font=f, fill=MUTE)
        img.save(path)
        return img

    def _gradient(self, W, H, top, bottom):
        base = Image.new("RGB", (W, H), top)
        # diagonal: blend along x+y
        grad = Image.new("L", (W, H))
        gd = grad.load()
        for y in range(0, H, 2):
            for x in range(0, W, 2):
                t = (x / W * 0.4 + y / H * 0.6)
                v = int(t * 255)
                gd[x, y] = v
                if x + 1 < W:
                    gd[x + 1, y] = v
                if y + 1 < H:
                    gd[x, y + 1] = v
                    if x + 1 < W:
                        gd[x + 1, y + 1] = v
        return Image.composite(Image.new("RGB", (W, H), bottom), base, grad)

    def _draw(self, d, img, el):
        kind = el[0]
        if kind == "rect":
            _, x, y, w, h, c, r = el
            box = [self._I(x), self._I(y), self._I(x + w), self._I(y + h)]
            if r and r > 0:
                d.rounded_rectangle(box, radius=self._I(r), fill=c)
            else:
                d.rectangle(box, fill=c)
        elif kind == "line":
            _, x1, y1, x2, y2, c, w = el
            d.line([self._I(x1), self._I(y1), self._I(x2), self._I(y2)], fill=c, width=max(1, int(w)))
        elif kind == "arrow":
            _, x1, y1, x2, y2, c, w = el
            self._arrow(d, x1, y1, x2, y2, c, w)
        elif kind == "image":
            _, x, y, w, h, im = el
            im2 = im.convert("RGBA").resize((max(1, self._I(w)), max(1, self._I(h))))
            img.paste(im2, (self._I(x), self._I(y)), im2)
        elif kind == "text":
            _, x, y, w, h, t, size, c, bold, italic, align, ls = el
            f = lambda cj: _font(self._pt(size), bold, italic, cj)
            maxw = self._I(w)
            lh = int(self._pt(size) * 1.18 * ls)
            yy = self._I(y)
            for raw in t.split("\n"):
                ff = f(_has_cjk(raw))
                for line in self._wrap(d, raw, ff, maxw):
                    tw = d.textlength(line, font=ff)
                    xx = self._I(x) + (maxw - tw) // 2 if align == "c" else \
                         self._I(x) + maxw - tw if align == "r" else self._I(x)
                    d.text((xx, yy), line, font=ff, fill=c)
                    yy += lh
        elif kind == "bullets":
            _, x, y, w, h, items, size, c, gap = el
            maxw = self._I(w)
            lh = int(self._pt(size) * 1.32)
            yy = self._I(y)
            for it in items:
                cj = _has_cjk(it)
                ff = _font(self._pt(size), False, False, cj)
                # dot marker in accent
                dot_r = max(2, self._pt(size) // 7)
                cx = self._I(x) + dot_r + 1
                cyy = yy + int(self._pt(size) * 0.55)
                d.ellipse([cx - dot_r, cyy - dot_r, cx + dot_r, cyy + dot_r], fill=self.accent)
                tx = self._I(x) + dot_r * 2 + self._pt(size) // 3
                for j, line in enumerate(self._wrap(d, it, ff, maxw - (tx - self._I(x)))):
                    d.text((tx, yy), line, font=ff, fill=c)
                    yy += lh
                yy += int(gap * self._PX / 72.0)

    def _arrow(self, d, x1, y1, x2, y2, c, w):
        X1, Y1, X2, Y2 = self._I(x1), self._I(y1), self._I(x2), self._I(y2)
        d.line([X1, Y1, X2, Y2], fill=c, width=max(2, int(w)))
        ang = math.atan2(Y2 - Y1, X2 - X1)
        L = self._I(0.16)
        for da in (math.radians(150), math.radians(-150)):
            xa = X2 + L * math.cos(ang + da)
            ya = Y2 + L * math.sin(ang + da)
            d.line([X2, Y2, xa, ya], fill=c, width=max(2, int(w)))

    def _contact(self, previews, path):
        imgs = [Image.open(p) for p in previews]
        cols = 4
        rows = math.ceil(len(imgs) / cols)
        tw, th = 320, 180
        pad = 10
        sheet = Image.new("RGB", (cols * tw + pad * (cols + 1), rows * th + pad * (rows + 1)), (225, 228, 232))
        for i, im in enumerate(imgs):
            r, c = divmod(i, cols)
            sheet.paste(im.resize((tw, th)), (pad + c * (tw + pad), pad + r * (th + pad)))
        sheet.save(path)

    def _pdf(self, name, outdir, path):
        old = self._PX
        self._PX = 170
        _fc.clear()
        try:
            pages = [self._png(s, i + 1, os.path.join(outdir, f"_tmp_{i}.png"))
                     for i, s in enumerate(self.slides)]
            pages = [p.convert("RGB") for p in pages]
            pages[0].save(path, "PDF", save_all=True, append_images=pages[1:], resolution=170.0)
            for i in range(len(self.slides)):
                os.remove(os.path.join(outdir, f"_tmp_{i}.png"))
        finally:
            self._PX = old
            _fc.clear()

    # ---- PPTX ----
    def _pptx(self, path):
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE

        def rgb(c):
            return RGBColor(*c)

        def run(p, text, size, color, bold, italic=False):
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = rgb(color)
            r.font.name = PPTX_JP_FONT if _has_cjk(text) else PPTX_LAT_FONT

        prs = Presentation()
        prs.slide_width = Inches(SW_IN)
        prs.slide_height = Inches(SH_IN)
        blank = prs.slide_layouts[6]
        for s in self.slides:
            sl = prs.slides.add_slide(blank)
            sl.background.fill.solid()
            sl.background.fill.fore_color.rgb = rgb(s.grad[0] if s.grad else s.bg)
            for el in s.el:
                k = el[0]
                if k == "rect":
                    _, x, y, w, h, c, r = el
                    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if r and r > 0 else MSO_SHAPE.RECTANGLE
                    sp = sl.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
                    sp.fill.solid(); sp.fill.fore_color.rgb = rgb(c)
                    sp.line.fill.background(); sp.shadow.inherit = False
                elif k in ("line", "arrow"):
                    _, x1, y1, x2, y2, c, w = el
                    cn = sl.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
                    cn.line.color.rgb = rgb(c); cn.line.width = Pt(w)
                elif k == "image":
                    _, x, y, w, h, im = el
                    buf = BytesIO(); im.save(buf, "PNG"); buf.seek(0)
                    sl.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))
                elif k == "text":
                    _, x, y, w, h, t, size, c, bold, italic, align, ls = el
                    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                    tf = tb.text_frame; tf.word_wrap = True
                    al = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
                    for i, line in enumerate(t.split("\n")):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.alignment = al; p.line_spacing = ls
                        run(p, line, size, c, bold, italic)
                elif k == "bullets":
                    _, x, y, w, h, items, size, c, gap = el
                    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                    tf = tb.text_frame; tf.word_wrap = True
                    for i, it in enumerate(items):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.space_after = Pt(gap); p.line_spacing = 1.2
                        run(p, "•  " + it, size, c, False)
        prs.save(path)
