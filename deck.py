#!/usr/bin/env python3
"""
Single source of truth for the 20-slide deck.
Generates BOTH the .pptx and PNG previews from one spec so they always match.
"""
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- palette (RGB tuples) ----------
NAVY=(0x0F,0x2A,0x43); SLATE=(0x33,0x44,0x55); ACCENT=(0xC0,0x8A,0x2B)
LIGHT=(0xF2,0xF4,0xF7); WHITE=(0xFF,0xFF,0xFF); GREY=(0x6B,0x76,0x82)
PANEL=(0xE8,0xEC,0xF1); FAINT=(0x1C,0x3A,0x56); DARKFOOT=(0x90,0x9C,0xAA)

SW_IN, SH_IN = 13.333, 7.5
PAGE_TOTAL=20
KEYBAND_Y, KEYBAND_H = 6.32, 0.62

# ================= SPEC BUILDERS =================
def slide(bg=WHITE): return {"bg":bg, "el":[]}
def rect(s,x,y,w,h,c): s["el"].append(("rect",x,y,w,h,c))
def text(s,x,y,w,t,size=18,color=SLATE,bold=False,italic=False,align="l",ls=1.0,h=0.8):
    s["el"].append(("text",x,y,w,h,t,size,color,bold,italic,align,ls))
def bullets(s,x,y,w,items,size=18,color=SLATE,gap=6,bold_lead=False,h=5.0):
    s["el"].append(("bullets",x,y,w,h,items,size,color,gap,bold_lead))

def header(s,kicker,title):
    rect(s,0,0,SW_IN,1.25,NAVY); rect(s,0,1.25,SW_IN,0.06,ACCENT)
    text(s,0.6,0.20,11.5,kicker,13,ACCENT,bold=True)
    text(s,0.6,0.52,12.1,title,26,WHITE,bold=True)

def keyband(s,txt):
    """Bottom key-line band that mirrors the top header band for balance."""
    rect(s,0,KEYBAND_Y,SW_IN,KEYBAND_H,LIGHT)
    rect(s,0,KEYBAND_Y,0.14,KEYBAND_H,ACCENT)
    text(s,0.62,KEYBAND_Y+0.15,12.1,txt,15,NAVY,italic=True,h=0.5)

def footer(s,n,dark=False):
    col=DARKFOOT if dark else GREY
    rect(s,0.6,7.2,0.45,0.035,ACCENT)
    text(s,11.4,7.02,1.33,f"{n:02d} / {PAGE_TOTAL}",11,col,align="r",h=0.3)

def divider(num,title,subtitle):
    s=slide(NAVY)
    text(s,7.7,0.7,5.4,f"{num:02d}",210,FAINT,bold=True,align="c",h=6.0)
    text(s,0.9,2.55,8.0,f"SECTION {num} OF 4",13,ACCENT,bold=True)
    rect(s,0.9,2.95,3.0,0.06,ACCENT)
    text(s,0.9,3.18,8.2,title,33,WHITE,bold=True,h=1.6)
    text(s,0.9,4.35,7.8,subtitle,18,LIGHT,italic=True,ls=1.15,h=1.2)
    return s

SLIDES=[]
def add(s): SLIDES.append(s); return s

# 1 — TITLE
s=add(slide(NAVY))
rect(s,0,0,SW_IN,0.16,ACCENT); rect(s,0,SH_IN-0.16,SW_IN,0.16,ACCENT)
text(s,1.0,2.5,11.3,"TRANSNATIONAL CRIME & THE LIMITS OF INTERNATIONAL JUSTICE",16,ACCENT,bold=True,align="c")
text(s,1.0,3.0,11.3,"What the Epstein Case Reveals",40,WHITE,bold=True,align="c",h=1.1)
rect(s,(SW_IN-3)/2,4.12,3,0.06,ACCENT)
text(s,1.0,4.4,11.3,"A case study in cross-border exploitation, jurisdiction, and accountability",18,LIGHT,italic=True,align="c")
text(s,1.0,6.55,11.3,"Group Presentation  •  May 2026",13,DARKFOOT,align="c")

# 2 — AGENDA
s=add(slide()); header(s,"OVERVIEW","Our Central Question & Roadmap")
rect(s,0.6,1.65,12.13,1.2,LIGHT); rect(s,0.6,1.65,0.14,1.2,ACCENT)
text(s,0.85,1.85,11.6,"Why is transnational sexual exploitation so hard to prosecute — and what do international law and cross-border cooperation need in order to work?",19,NAVY,bold=True,italic=True,ls=1.1,h=1.0)
bullets(s,0.8,3.35,11.8,[
 "Part 1 — Framing & the international scope of the case",
 "Part 2 — The international legal framework (Palermo Protocol, UNTOC, jurisdiction & extradition)",
 "Part 3 — Wealth, mobility, and the problem of \"elite impunity\"",
 "Part 4 — Accountability, transparency & lessons (incl. a Japan comparison)",
],19,SLATE,gap=20)
keyband(s,"This presentation centres victims and legal systems — not sensational detail.")

# 3 — DIVIDER P1
add(divider(1,"Framing & the\nInternational Scope","Who Epstein was — and why this case crosses borders"))

# 4 — WHO WAS EPSTEIN
s=add(slide()); header(s,"PART 1  •  BACKGROUND","Who Was Jeffrey Epstein?")
bullets(s,0.7,1.95,7.4,[
 "A US financier who built extraordinary wealth and a network of powerful international contacts.",
 "2005–06: A Palm Beach (Florida) police investigation into abuse of minors.",
 "2008: A controversial state plea deal; he registered as a sex offender.",
 "2019: Arrested on federal sex-trafficking charges in New York.",
 "Died in custody in August 2019 (ruled a suicide) — ending the criminal case against him.",
],18,gap=26,bold_lead=True)
rect(s,8.45,1.75,4.28,4.35,LIGHT); rect(s,8.45,1.75,0.14,4.35,ACCENT)
text(s,8.75,2.0,3.8,"WHY HE MATTERS HERE",14,ACCENT,bold=True)
text(s,8.75,2.55,3.8,"His death meant accountability had to be pursued through others — and across borders. The case became a test of how legal systems cooperate.",16,SLATE,italic=True,ls=1.2,h=3.3)
keyband(s,"The crime was global; the criminal case ended at one nation's border.")

# 5 — TIMELINE
s=add(slide()); header(s,"PART 1  •  CHRONOLOGY","A Case That Spanned Two Decades")
stages=[("2008","Florida state plea deal — widely criticised as lenient."),
        ("2019","Federal trafficking charges in New York; death in custody."),
        ("2021","Ghislaine Maxwell convicted of sex-trafficking conspiracy."),
        ("2025","Document-release wave renews scrutiny and debate.")]
bw=2.85; gp=0.18; x0=0.7; y0=2.55
rect(s,x0+bw/2,y0+0.41,3*(bw+gp),0.05,ACCENT)   # connecting line behind boxes
for i,(yr,cap) in enumerate(stages):
    x=x0+i*(bw+gp)
    rect(s,x,y0,bw,0.85,NAVY)
    text(s,x,y0+0.17,bw,yr,28,WHITE,bold=True,align="c")
    rect(s,x+bw/2-0.03,y0+0.85,0.06,0.33,ACCENT)
    text(s,x,y0+1.32,bw,cap,14,SLATE,ls=1.15,h=2.0)
keyband(s,"From a local plea bargain to a global question about justice, wealth and cross-border accountability.")

# 6 — WHY INTERNATIONAL
s=add(slide()); header(s,"PART 1  •  SCOPE","Why This Is an International Story")
bullets(s,0.7,1.95,7.4,[
 "Recruitment: victims, including minors and foreign nationals, moved across borders.",
 "Finance: wealth held through offshore structures, obscuring ownership.",
 "Property: residences across the US, the US Virgin Islands and Paris.",
 "Identity: an expired passport bore his photo, a false name, and listed Saudi Arabia as residence.",
 "Network: associates spanning several countries and high office.",
],18,gap=26,bold_lead=True)
rect(s,8.45,1.75,4.28,4.35,LIGHT); rect(s,8.45,1.75,0.14,4.35,ACCENT)
text(s,8.75,2.0,3.8,"THE CORE TENSION",14,ACCENT,bold=True)
text(s,8.75,2.55,3.8,"The network was global, but criminal enforcement is national. That gap is the heart of our presentation.",16,SLATE,italic=True,ls=1.2,h=3.3)
keyband(s,"A global network meets national enforcement — and falls through the gaps.")

# 7 — DIVIDER P2
add(divider(2,"The International\nLegal Framework","How international law addresses crime that crosses borders"))

# 8 — TRANSNATIONAL CRIME DEFINED
s=add(slide()); header(s,"PART 2  •  CONCEPTS","What Makes a Crime \"Transnational\"?")
bullets(s,0.7,2.1,11.9,[
 "Committed in more than one state — or in one state but planned, directed or with effects in another.",
 "Trafficking in persons is a paradigm case: it inherently moves people, money and evidence across borders.",
 "Prosecuting it therefore requires more than one country's police, courts and laws to act together.",
 "This is exactly where domestic systems, built for domestic crime, come under strain.",
],19,gap=30)
keyband(s,"The question is institutional, not just moral: can sovereign systems coordinate fast enough?")

# 9 — PALERMO
s=add(slide()); header(s,"PART 2  •  INSTRUMENTS","The Palermo Protocol (2000)")
bullets(s,0.7,2.05,11.9,[
 "Full name: UN Protocol to Prevent, Suppress and Punish Trafficking in Persons.",
 "First globally agreed definition of human trafficking.",
 "Three elements: the act, the means, and the purpose of exploitation.",
 "For children, the \"means\" element is not required — any recruitment for exploitation qualifies.",
 "Obliges states parties to criminalise trafficking in their domestic law.",
],19,gap=28,bold_lead=True)
keyband(s,"It gives 190+ states a shared vocabulary — a precondition for cooperation.")

# 10 — UNTOC
s=add(slide()); header(s,"PART 2  •  INSTRUMENTS","UNTOC & Cross-Border Cooperation")
bullets(s,0.7,2.05,11.9,[
 "Parent treaty: the UN Convention against Transnational Organized Crime.",
 "The Palermo Protocol supplements it; the two operate together.",
 "Provides a framework for extradition, mutual legal assistance and joint investigations.",
 "Can serve as a legal basis for cooperation even where no bilateral treaty exists.",
 "Encourages asset confiscation and protection of victims and witnesses.",
],19,gap=28,bold_lead=True)
keyband(s,"In short: the machinery for international cooperation already exists on paper.")

# 11 — JURISDICTION / EXTRADITION
s=add(slide()); header(s,"PART 2  •  ENFORCEMENT","Jurisdiction, Extradition & MLATs")
bullets(s,0.7,2.05,11.9,[
 "Jurisdiction: states may claim it by territory or by nationality of offender/victim.",
 "Concurrent claims: several states may have a right to prosecute the same conduct.",
 "Extradition: typically needs a treaty plus \"dual criminality\" (an offence in both states).",
 "MLATs: Mutual Legal Assistance Treaties move evidence, testimony and records across borders.",
 "The friction: requests are slow, can be refused, and stall when suspects are wealthy or well-connected.",
],19,gap=28,bold_lead=True)
keyband(s,"The tools exist — but speed and political will decide whether they work.")

# 12 — FOREIGN NATIONALS
s=add(slide()); header(s,"PART 2  •  ACTORS","Foreign Nationals & Quasi-Diplomatic Figures")
bullets(s,0.7,2.1,11.9,[
 "Ghislaine Maxwell (UK): convicted in 2021 of sex-trafficking conspiracy in a US court.",
 "Prince Andrew: settled a US civil claim in 2022 and later stepped back from royal titles.",
 "High status confers no criminal immunity — but it can create practical and diplomatic hurdles.",
 "Foreign nationality complicates arrest, extradition and the gathering of evidence abroad.",
],19,gap=30,bold_lead=True)
keyband(s,"Stick to what is adjudicated: charges, a conviction, a settled civil case — not speculation.")

# 13 — DIVIDER P3
add(divider(3,"Wealth, Mobility\n& Impunity","How resources can delay — though not always defeat — accountability"))

# 14 — OFFSHORE WEALTH
s=add(slide()); header(s,"PART 3  •  MECHANICS","Offshore Wealth & Global Mobility")
bullets(s,0.7,2.05,11.9,[
 "Offshore structures and trusts can obscure who really owns assets.",
 "Property in several jurisdictions provides places to live, store value and move between.",
 "Private travel reduces the friction and scrutiny of ordinary border crossings.",
 "Multiple identity documents — including the false-name passport — illustrate evasion capacity.",
 "Together these convert wealth into mobility, and mobility into delay.",
],19,gap=28)
keyband(s,"None of this is unique to one man — it is how illicit elite wealth tends to behave.")

# 15 — ELITE IMPUNITY (compare)
s=add(slide()); header(s,"PART 3  •  COMPARISON","\"Elite Impunity\": 2008 vs 2019")
rect(s,0.7,1.8,5.85,4.25,PANEL); rect(s,0.7,1.8,0.14,4.25,ACCENT)
text(s,1.05,2.05,5.3,"2008 — FLORIDA",16,ACCENT,bold=True)
bullets(s,1.05,2.6,5.2,[
 "State-level plea deal",
 "Served ~13 months, with work release",
 "Non-prosecution agreement shielded others",
 "Victims said they were not consulted",
],15,SLATE,gap=16)
rect(s,6.8,1.8,5.85,4.25,NAVY); rect(s,6.8,1.8,0.14,4.25,ACCENT)
text(s,7.15,2.05,5.3,"2019 — FEDERAL (NY)",16,ACCENT,bold=True)
bullets(s,7.15,2.6,5.2,[
 "Federal sex-trafficking charges",
 "Held without bail",
 "Treated as serious organised wrongdoing",
 "Case ended only with his death in custody",
],15,LIGHT,gap=16)
keyband(s,"The same conduct, two very different responses — a decade apart.")

# 16 — COMPARATIVE LAW
s=add(slide()); header(s,"PART 3  •  COMPARATIVE LAW","Would Other Systems Have Acted Differently?")
bullets(s,0.7,2.1,11.9,[
 "Adversarial vs inquisitorial: France's investigating magistrates can drive cases differently from US prosecutors.",
 "Victim participation: some systems give victims a formal role earlier in proceedings.",
 "Limitation periods: how long after the abuse can charges still be brought?",
 "Plea bargaining: less central in many civil-law systems, changing the incentive to settle quietly.",
],19,gap=30,bold_lead=True)
keyband(s,"Central question: would the UK, France or Japan have acted faster — or slower?")

# 17 — DIVIDER P4
add(divider(4,"Accountability,\nTransparency & Lessons","Designing cross-border justice that actually protects victims"))

# 18 — 2025 TRANSPARENCY
s=add(slide()); header(s,"PART 4  •  TRANSPARENCY","The 2025 Document-Release Wave")
bullets(s,0.7,2.1,11.9,[
 "US transparency legislation drove a large release of case-related documents.",
 "Intended to answer public demand for accountability and openness.",
 "Survivors' concern: the disclosures made it easy to identify those abused — but not those who may have been involved.",
 "A real tension emerges between public transparency and the privacy and safety of victims.",
],19,gap=30,bold_lead=True)
keyband(s,"Transparency is not automatically just — its design determines who it protects.")

# 19 — REFORM + JAPAN
s=add(slide()); header(s,"PART 4  •  LESSONS","Reform Proposals & a Japan Comparison")
bullets(s,0.7,1.95,7.4,[
 "Faster, better-resourced MLAT and extradition channels.",
 "Victim-centred disclosure rules that shield identities by default.",
 "Stronger cross-border asset tracing and recovery.",
 "Harmonised definitions so cooperation isn't blocked by mismatched laws.",
],18,gap=26,bold_lead=False)
rect(s,8.45,1.75,4.28,4.35,LIGHT); rect(s,8.45,1.75,0.14,4.35,ACCENT)
text(s,8.75,2.0,3.8,"JAPAN ANGLE",14,ACCENT,bold=True)
text(s,8.75,2.55,3.8,"Japan has an anti-trafficking action framework and is assessed annually in the US State Department's TIP Report — a useful benchmark for how a non-Western system measures up on cross-border cooperation.",15,SLATE,italic=True,ls=1.22,h=4.0)
keyband(s,"Better justice across borders is mostly an engineering problem, not a gap in principle.")

# 20 — CONCLUSION + sources
s=add(slide()); header(s,"CONCLUSION","Four Takeaways")
items=[("01","Crime crosses borders — justice often doesn't.","Networks are global; enforcement is national."),
 ("02","The legal tools already exist.","Palermo & UNTOC provide the framework; coordination is the gap."),
 ("03","Wealth buys time, not immunity.","Mobility delayed — but did not finally prevent — accountability."),
 ("04","Transparency must protect victims.","Disclosure should expose enablers, not endanger survivors.")]
y=1.6
for num,head_,sub in items:
    rect(s,0.7,y,0.8,0.9,ACCENT)
    text(s,0.7,y+0.18,0.8,num,24,WHITE,bold=True,align="c")
    text(s,1.65,y+0.06,10.9,head_,18,NAVY,bold=True)
    text(s,1.65,y+0.46,10.9,sub,14,SLATE,italic=True)
    y+=1.02
rect(s,0.6,5.95,12.13,0.95,LIGHT); rect(s,0.6,5.95,0.14,0.95,ACCENT)
text(s,0.85,6.08,11.6,"SELECTED SOURCES & INSTRUMENTS",11,ACCENT,bold=True)
text(s,0.85,6.4,11.7,"Palermo Protocol (2000) · UN Convention against Transnational Organized Crime · US DOJ filings (US v. Epstein 2019; US v. Maxwell 2021) · US State Dept. Trafficking in Persons Reports · contemporary reporting on 2025 disclosures.",11,SLATE,ls=1.15,h=0.55)

# ---- footers / page numbers on every slide except the title ----
for i,sp in enumerate(SLIDES):
    if i==0: continue
    footer(sp,i+1,dark=(sp["bg"]==NAVY))

# ================= PPTX BACKEND =================
def to_rgb(c): return RGBColor(*c)
def build_pptx(path):
    prs=Presentation(); prs.slide_width=Inches(SW_IN); prs.slide_height=Inches(SH_IN)
    blank=prs.slide_layouts[6]
    for spec in SLIDES:
        sl=prs.slides.add_slide(blank)
        sl.background.fill.solid(); sl.background.fill.fore_color.rgb=to_rgb(spec["bg"])
        for el in spec["el"]:
            kind=el[0]
            if kind=="rect":
                _,x,y,w,h,c=el
                sp=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
                sp.fill.solid(); sp.fill.fore_color.rgb=to_rgb(c)
                sp.line.fill.background(); sp.shadow.inherit=False
            elif kind=="text":
                _,x,y,w,h,t,size,color,bold,italic,align,ls=el
                tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
                tf=tb.text_frame; tf.word_wrap=True
                al={"l":PP_ALIGN.LEFT,"c":PP_ALIGN.CENTER,"r":PP_ALIGN.RIGHT}[align]
                for i,line in enumerate(t.split("\n")):
                    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
                    p.alignment=al; p.line_spacing=ls
                    r=p.add_run(); r.text=line
                    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
                    r.font.color.rgb=to_rgb(color); r.font.name="Times New Roman"
            elif kind=="bullets":
                _,x,y,w,h,items_,size,color,gap,bold_lead=el
                tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
                tf=tb.text_frame; tf.word_wrap=True
                for i,item in enumerate(items_):
                    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
                    p.space_after=Pt(gap); p.line_spacing=1.05
                    if bold_lead and ": " in item:
                        lead,rest=item.split(": ",1)
                        r1=p.add_run(); r1.text="•  "+lead+": "
                        r1.font.size=Pt(size); r1.font.bold=True
                        r1.font.color.rgb=to_rgb(NAVY); r1.font.name="Times New Roman"
                        r2=p.add_run(); r2.text=rest
                        r2.font.size=Pt(size); r2.font.color.rgb=to_rgb(color); r2.font.name="Times New Roman"
                    else:
                        r=p.add_run(); r.text="•  "+item
                        r.font.size=Pt(size); r.font.color.rgb=to_rgb(color); r.font.name="Times New Roman"
    prs.save(path)
    return len(SLIDES)

# ================= PNG BACKEND =================
PX=96
def I(v): return int(round(v*PX))
def pt(p): return int(round(p*PX/72.0))
SERIF="/usr/share/fonts/truetype/liberation/LiberationSerif-Regular.ttf"
SERIFB="/usr/share/fonts/truetype/liberation/LiberationSerif-Bold.ttf"
SERIFI="/usr/share/fonts/truetype/liberation/LiberationSerif-Italic.ttf"
SERIFBI="/usr/share/fonts/truetype/liberation/LiberationSerif-BoldItalic.ttf"
_fc={}
def font(sz,bold=False,italic=False):
    path=SERIFBI if(bold and italic) else SERIFB if bold else SERIFI if italic else SERIF
    k=(path,sz)
    if k not in _fc: _fc[k]=ImageFont.truetype(path,sz)
    return _fc[k]
def wrap(d,t,f,maxw):
    out=[]; cur=""
    for w in t.split():
        cand=(cur+" "+w).strip()
        if d.textlength(cand,font=f)<=maxw: cur=cand
        else:
            if cur: out.append(cur)
            cur=w
    if cur: out.append(cur)
    return out or [""]

def render_png(spec,path):
    img=Image.new("RGB",(I(SW_IN),I(SH_IN)),spec["bg"]); d=ImageDraw.Draw(img)
    for el in spec["el"]:
        kind=el[0]
        if kind=="rect":
            _,x,y,w,h,c=el; d.rectangle([I(x),I(y),I(x+w),I(y+h)],fill=c)
        elif kind=="text":
            _,x,y,w,h,t,size,color,bold,italic,align,ls=el
            f=font(pt(size),bold,italic); maxw=I(w); lh=int(pt(size)*1.2*ls); yy=I(y)
            for raw in t.split("\n"):
                for line in wrap(d,raw,f,maxw):
                    tw=d.textlength(line,font=f)
                    xx=I(x)+(maxw-tw)//2 if align=="c" else I(x)+maxw-tw if align=="r" else I(x)
                    d.text((xx,yy),line,font=f,fill=color); yy+=lh
        elif kind=="bullets":
            _,x,y,w,h,items_,size,color,gap,bold_lead=el
            f=font(pt(size)); fb=font(pt(size),bold=True); maxw=I(w); lh=int(pt(size)*1.18); yy=I(y)
            indent=I(x)+d.textlength("•  ",font=f)
            for it in items_:
                if bold_lead and ": " in it:
                    lead,rest=it.split(": ",1)
                    words=[(w_,fb) for w_ in ("•  "+lead+": ").split(" ") if w_!=""]
                    words+=[(w_,f) for w_ in rest.split(" ") if w_!=""]
                else:
                    words=[(w_,f) for w_ in ("•  "+it).split(" ") if w_!=""]
                line=[]; lw=0; firstline=True
                def flush(line,yy,first):
                    xpos=I(x) if first else indent
                    for wd,ff in line:
                        d.text((xpos,yy),wd+" ",font=ff,fill=(NAVY if ff is fb else color))
                        xpos+=d.textlength(wd+" ",font=ff)
                for wd,ff in words:
                    ww=d.textlength(wd+" ",font=ff)
                    avail=maxw-(0 if firstline else (indent-I(x)))
                    if lw+ww>avail and line:
                        flush(line,yy,firstline); yy+=lh; line=[]; lw=0; firstline=False
                    line.append((wd,ff)); lw+=ww
                if line: flush(line,yy,firstline); yy+=lh
                yy+=int(gap*PX/72.0)
    img.save(path); return img

if __name__=="__main__":
    n=build_pptx("/home/user/Koki/Epstein_Transnational_Justice.pptx")
    imgs=[render_png(s,f"/home/user/Koki/preview_slide_{i+1}.png") for i,s in enumerate(SLIDES)]
    cols,rows=4,5; tw,th=308,173; pad=12
    sheet=Image.new("RGB",(cols*tw+pad*(cols+1),rows*th+pad*(rows+1)),(210,214,219))
    for idx,im in enumerate(imgs):
        r,c=divmod(idx,cols)
        sheet.paste(im.resize((tw,th)),(pad+c*(tw+pad),pad+r*(th+pad)))
    sheet.save("/home/user/Koki/preview_all.png")
    print(f"built {n} slides + {len(imgs)} previews + contact sheet")
