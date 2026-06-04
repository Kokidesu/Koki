#!/usr/bin/env python3
"""Presentation God — a reusable, beautiful deck engine.

One spec → .pptx + high-res .pdf + PNG previews + a contact sheet, so the
editable deck and the rendered previews can never drift apart.

Harvested and generalised from the hand-tuned `transnational-justice/deck.py`:
a clean consulting aesthetic (navy + gold accent), section dividers, key-line
bands, bold-lead bullets, line-art icons — now driven by data and bilingual
(Japanese / English) with automatic CJK-aware font selection.

Usage:
    from engine import Deck, THEMES
    d = Deck(theme="consulting")
    d.title("KICKER", "Big Title", "subtitle", "footer note")
    d.section(1, 4, "Section Title", "one-line description")
    d.bullets("PART 1 · TOPIC", "Slide title", ["point one", "Lead: detail"],
              key="bottom key-line", bold_lead=True)
    d.takeaways("CONCLUSION", "Takeaways", [("01","head","sub"), ...], sources="...")
    d.build("/out/dir", "my_deck")
"""
from __future__ import annotations

import math
import os
from dataclasses import dataclass, field
from io import BytesIO
from typing import Any

from PIL import Image, ImageDraw, ImageFont

# ---------- canvas ----------
SW_IN, SH_IN = 13.333, 7.5
KEYBAND_Y, KEYBAND_H = 6.32, 0.62
FONT_SCALE = 1.15


# ---------- themes ----------
@dataclass
class Theme:
    name: str
    navy: tuple        # primary dark
    accent: tuple      # accent / gold
    slate: tuple       # body text
    light: tuple       # light panel
    white: tuple
    grey: tuple
    panel: tuple
    faint: tuple       # huge ghost numerals on dividers
    darkfoot: tuple


THEMES = {
    "consulting": Theme(
        name="consulting",
        navy=(0x0F, 0x2A, 0x43), accent=(0xC0, 0x8A, 0x2B), slate=(0x33, 0x44, 0x55),
        light=(0xF2, 0xF4, 0xF7), white=(0xFF, 0xFF, 0xFF), grey=(0x6B, 0x76, 0x82),
        panel=(0xE8, 0xEC, 0xF1), faint=(0x1C, 0x3A, 0x56), darkfoot=(0x90, 0x9C, 0xAA),
    ),
    "midnight": Theme(
        name="midnight",
        navy=(0x10, 0x14, 0x2B), accent=(0x6E, 0xA8, 0xFE), slate=(0x2A, 0x30, 0x45),
        light=(0xEE, 0xF1, 0xF8), white=(0xFF, 0xFF, 0xFF), grey=(0x68, 0x70, 0x85),
        panel=(0xE5, 0xEA, 0xF4), faint=(0x22, 0x2A, 0x4A), darkfoot=(0x8A, 0x93, 0xAD),
    ),
    "forest": Theme(
        name="forest",
        navy=(0x12, 0x33, 0x2A), accent=(0xCB, 0x8A, 0x3E), slate=(0x2E, 0x3E, 0x39),
        light=(0xEF, 0xF4, 0xF1), white=(0xFF, 0xFF, 0xFF), grey=(0x66, 0x72, 0x6D),
        panel=(0xE2, 0xEC, 0xE7), faint=(0x1E, 0x42, 0x37), darkfoot=(0x8C, 0x9A, 0x94),
    ),
}


# ---------- fonts (CJK-aware) ----------
# English: Liberation Serif locally (metric-compatible with Times New Roman),
#          and "Times New Roman" written into the .pptx so PowerPoint shows it.
# Japanese: Noto Serif CJK JP locally (an elegant mincho serif that pairs with
#          Times New Roman), and 游明朝 / "Yu Mincho" written into the .pptx.
_LATIN = {
    "regular": "/usr/share/fonts/truetype/liberation/LiberationSerif-Regular.ttf",
    "bold": "/usr/share/fonts/truetype/liberation/LiberationSerif-Bold.ttf",
    "italic": "/usr/share/fonts/truetype/liberation/LiberationSerif-Italic.ttf",
    "bolditalic": "/usr/share/fonts/truetype/liberation/LiberationSerif-BoldItalic.ttf",
}
# (path, index) for the .ttc collections; JP is index 0 in Noto's CJK files.
_CJK_REG = ("/usr/share/fonts/opentype/noto/NotoSerifCJK-Regular.ttc", 0)
_CJK_BOLD = ("/usr/share/fonts/opentype/noto/NotoSerifCJK-Bold.ttc", 0)
_CJK_FALLBACK = ("/usr/share/fonts/opentype/ipafont-gothic/ipagp.ttf", 0)
if not os.path.exists(_CJK_REG[0]):
    _CJK_REG = _CJK_BOLD = _CJK_FALLBACK

_CJK_RE = __import__("re").compile(r"[぀-ヿ぀-ゟ゠-ヿ一-鿿０-ｚ]")

# pptx font names that resolve beautifully on both Windows and Mac PowerPoint.
PPTX_LATIN_FONT = "Times New Roman"
PPTX_CJK_FONT = "Yu Mincho"  # 游明朝 — elegant, ships with Windows & macOS


def _has_cjk(text: str) -> bool:
    return bool(_CJK_RE.search(text))


_fontcache: dict[tuple, ImageFont.FreeTypeFont] = {}


def _pil_font(size: int, bold: bool, italic: bool, cjk: bool):
    if cjk:
        path, index = _CJK_BOLD if bold else _CJK_REG
    else:
        key = "bolditalic" if (bold and italic) else "bold" if bold else "italic" if italic else "regular"
        path, index = _LATIN[key], 0
    ck = (path, index, size)
    if ck not in _fontcache:
        _fontcache[ck] = ImageFont.truetype(path, size, index=index)
    return _fontcache[ck]


# ================= DECK MODEL =================
class Deck:
    def __init__(self, theme: str = "consulting", page_total: int | None = None):
        self.t = THEMES[theme]
        self.slides: list[dict[str, Any]] = []
        self.page_total = page_total  # set automatically at build if None
        self._icons: dict[str, Image.Image] = {}

    # ---- low-level spec builders ----
    def _slide(self, bg=None):
        s = {"bg": bg or self.t.white, "el": []}
        self.slides.append(s)
        return s

    @staticmethod
    def _rect(s, x, y, w, h, c):
        s["el"].append(("rect", x, y, w, h, c))

    @staticmethod
    def _text(s, x, y, w, t, size=18, color=None, bold=False, italic=False, align="l", ls=1.0, h=0.8):
        s["el"].append(("text", x, y, w, h, t, size, color, bold, italic, align, ls))

    @staticmethod
    def _bullets(s, x, y, w, items, size=18, color=None, gap=6, bold_lead=False, h=5.0):
        s["el"].append(("bullets", x, y, w, h, items, size, color, gap, bold_lead))

    @staticmethod
    def _image(s, x, y, w, h, name):
        s["el"].append(("image", x, y, w, h, name))

    def _header(self, s, kicker, title, icon=None):
        t = self.t
        self._rect(s, 0, 0, SW_IN, 1.25, t.navy)
        self._rect(s, 0, 1.25, SW_IN, 0.06, t.accent)
        self._text(s, 0.6, 0.20, 11.5, kicker, 13, t.accent, bold=True)
        self._text(s, 0.6, 0.56, 11.6, title, 25, t.white, bold=True, h=0.7)
        if icon:
            self._image(s, 12.08, 0.33, 0.6, 0.6, icon)

    def _keyband(self, s, txt):
        t = self.t
        self._rect(s, 0, KEYBAND_Y, SW_IN, KEYBAND_H, t.light)
        self._rect(s, 0, KEYBAND_Y, 0.14, KEYBAND_H, t.accent)
        self._text(s, 0.62, KEYBAND_Y + 0.15, 12.1, txt, 15, t.navy, italic=True, h=0.5)

    def _footer(self, s, n, dark):
        t = self.t
        col = t.darkfoot if dark else t.grey
        self._rect(s, 0.6, 7.2, 0.45, 0.035, t.accent)
        self._text(s, 11.4, 7.02, 1.33, f"{n:02d} / {self.page_total:02d}", 11, col, align="r", h=0.3)

    # ---- high-level slide types ----
    def title(self, kicker, title, subtitle="", footnote="", icon="globe"):
        t = self.t
        s = self._slide(t.navy)
        self._rect(s, 0, 0, SW_IN, 0.16, t.accent)
        self._rect(s, 0, SH_IN - 0.16, SW_IN, 0.16, t.accent)
        self._text(s, 1.0, 2.35, 11.3, kicker, 16, t.accent, bold=True, align="c")
        self._text(s, 1.0, 2.85, 11.3, title, 38, t.white, bold=True, align="c", h=1.4, ls=1.05)
        self._rect(s, (SW_IN - 3) / 2, 4.35, 3, 0.06, t.accent)
        if subtitle:
            self._text(s, 1.0, 4.62, 11.3, subtitle, 18, t.light, italic=True, align="c", h=1.0, ls=1.1)
        if icon:
            self._image(s, (SW_IN - 1.1) / 2, 5.5, 1.1, 1.1, icon)
        if footnote:
            self._text(s, 1.0, 6.85, 11.3, footnote, 12, t.darkfoot, align="c")
        return s

    def section(self, num, total, title, subtitle="", icon=None):
        t = self.t
        s = self._slide(t.navy)
        self._text(s, 7.7, 0.7, 5.4, f"{num:02d}", 200, t.faint, bold=True, align="c", h=6.0)
        self._text(s, 0.9, 2.55, 8.0, f"SECTION {num} OF {total}", 13, t.accent, bold=True)
        self._rect(s, 0.9, 2.95, 3.0, 0.06, t.accent)
        self._text(s, 0.9, 3.18, 8.4, title, 32, t.white, bold=True, h=1.6, ls=1.05)
        if subtitle:
            self._text(s, 0.9, 4.5, 7.8, subtitle, 18, t.light, italic=True, ls=1.15, h=1.2)
        if icon:
            self._image(s, 10.55, 4.6, 1.85, 1.85, icon)
        return s

    def bullets(self, kicker, title, items, key="", size=18, gap=None, bold_lead=False, icon=None):
        s = self._slide()
        self._header(s, kicker, title, icon)
        # Auto-balance: spread the list through the body band so it never bunches
        # at the top with dead space above the key-line.
        n = len(items)
        if gap is None:
            gap = {1: 60, 2: 48, 3: 40, 4: 30, 5: 22}.get(n, 16)
        y0 = {1: 3.1, 2: 2.7, 3: 2.35, 4: 2.1}.get(n, 1.95)
        self._bullets(s, 0.7, y0, 11.9, items, size, self.t.slate, gap=gap, bold_lead=bold_lead)
        if key:
            self._keyband(s, key)
        return s

    def split(self, kicker, title, items, side_title, side_body, key="", bold_lead=True, icon=None):
        """Left bullet list + a right highlight panel."""
        t = self.t
        s = self._slide()
        self._header(s, kicker, title, icon)
        self._bullets(s, 0.7, 1.95, 7.4, items, 18, t.slate, gap=24, bold_lead=bold_lead)
        self._rect(s, 8.45, 1.75, 4.28, 4.35, t.light)
        self._rect(s, 8.45, 1.75, 0.14, 4.35, t.accent)
        self._text(s, 8.75, 2.0, 3.8, side_title, 14, t.accent, bold=True)
        self._text(s, 8.75, 2.55, 3.8, side_body, 16, t.slate, italic=True, ls=1.2, h=3.3)
        if key:
            self._keyband(s, key)
        return s

    def compare(self, kicker, title, left_title, left_items, right_title, right_items, key=""):
        """Two-column comparison: light panel vs navy panel."""
        t = self.t
        s = self._slide()
        self._header(s, kicker, title)
        self._rect(s, 0.7, 1.8, 5.85, 4.25, t.panel)
        self._rect(s, 0.7, 1.8, 0.14, 4.25, t.accent)
        self._text(s, 1.05, 2.05, 5.3, left_title, 16, t.accent, bold=True)
        self._bullets(s, 1.05, 2.6, 5.2, left_items, 15, t.slate, gap=16)
        self._rect(s, 6.8, 1.8, 5.85, 4.25, t.navy)
        self._rect(s, 6.8, 1.8, 0.14, 4.25, t.accent)
        self._text(s, 7.15, 2.05, 5.3, right_title, 16, t.accent, bold=True)
        self._bullets(s, 7.15, 2.6, 5.2, right_items, 15, t.light, gap=16)
        if key:
            self._keyband(s, key)
        return s

    def takeaways(self, kicker, title, items, sources=""):
        """Numbered takeaways with an optional sources strip."""
        t = self.t
        s = self._slide()
        self._header(s, kicker, title)
        y = 1.6
        for num, head, sub in items:
            self._rect(s, 0.7, y, 0.8, 0.9, t.accent)
            self._text(s, 0.7, y + 0.18, 0.8, num, 24, t.white, bold=True, align="c")
            self._text(s, 1.65, y + 0.06, 10.9, head, 18, t.navy, bold=True, h=0.5)
            self._text(s, 1.65, y + 0.5, 10.9, sub, 14, t.slate, italic=True, h=0.45)
            y += 1.02
        if sources:
            self._rect(s, 0.6, 5.95, 12.13, 0.95, t.light)
            self._rect(s, 0.6, 5.95, 0.14, 0.95, t.accent)
            self._text(s, 0.85, 6.08, 11.6, "SOURCES / 出典", 11, t.accent, bold=True)
            self._text(s, 0.85, 6.4, 11.7, sources, 11, t.slate, ls=1.15, h=0.55)
        return s

    # ================= RENDER BACKENDS =================
    def _finalise(self):
        if self.page_total is None:
            self.page_total = len(self.slides)
        for i, sp in enumerate(self.slides):
            if i == 0:
                continue
            self._footer(sp, i + 1, dark=(sp["bg"] == self.t.navy))

    def build(self, outdir: str, name: str, pdf: bool = True, contact: bool = True) -> dict:
        os.makedirs(outdir, exist_ok=True)
        self._gen_icons()
        self._finalise()
        out = {}
        out["pptx"] = os.path.join(outdir, f"{name}.pptx")
        self._build_pptx(out["pptx"])
        previews = []
        for i, s in enumerate(self.slides):
            p = os.path.join(outdir, f"{name}_slide_{i + 1:02d}.png")
            self._render_png(s, p)
            previews.append(p)
        out["previews"] = previews
        if contact:
            out["contact"] = os.path.join(outdir, f"{name}_contact.png")
            self._contact_sheet(previews, out["contact"])
        if pdf:
            out["pdf"] = os.path.join(outdir, f"{name}.pdf")
            self._build_pdf(out["pdf"])
        out["n"] = len(self.slides)
        return out

    # ---- pptx ----
    def _build_pptx(self, path):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE

        def rgb(c):
            return RGBColor(*c)

        def style_run(r, size, color, bold, italic, text):
            r.text = text
            r.font.size = Pt(size * FONT_SCALE)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = rgb(color)
            r.font.name = PPTX_CJK_FONT if _has_cjk(text) else PPTX_LATIN_FONT

        prs = Presentation()
        prs.slide_width = Inches(SW_IN)
        prs.slide_height = Inches(SH_IN)
        blank = prs.slide_layouts[6]
        for spec in self.slides:
            sl = prs.slides.add_slide(blank)
            sl.background.fill.solid()
            sl.background.fill.fore_color.rgb = rgb(spec["bg"])
            for el in spec["el"]:
                kind = el[0]
                if kind == "rect":
                    _, x, y, w, h, c = el
                    sp = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
                    sp.fill.solid(); sp.fill.fore_color.rgb = rgb(c)
                    sp.line.fill.background(); sp.shadow.inherit = False
                elif kind == "image":
                    _, x, y, w, h, name = el
                    buf = BytesIO(); self._icons[name].save(buf, "PNG"); buf.seek(0)
                    sl.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))
                elif kind == "text":
                    _, x, y, w, h, t, size, color, bold, italic, align, ls = el
                    color = color or self.t.slate
                    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                    tf = tb.text_frame; tf.word_wrap = True
                    al = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
                    for i, line in enumerate(t.split("\n")):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.alignment = al; p.line_spacing = ls
                        style_run(p.add_run(), size, color, bold, italic, line)
                elif kind == "bullets":
                    _, x, y, w, h, items_, size, color, gap, bold_lead = el
                    color = color or self.t.slate
                    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                    tf = tb.text_frame; tf.word_wrap = True
                    for i, item in enumerate(items_):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.space_after = Pt(gap); p.line_spacing = 1.05
                        if bold_lead and ": " in item:
                            lead, rest = item.split(": ", 1)
                            style_run(p.add_run(), size, self.t.navy, True, False, "•  " + lead + ": ")
                            style_run(p.add_run(), size, color, False, False, rest)
                        else:
                            style_run(p.add_run(), size, color, False, False, "•  " + item)
        prs.save(path)

    # ---- PNG / PDF ----
    _PX = 96

    def _I(self, v):
        return int(round(v * self._PX))

    def _pt(self, p):
        return int(round(p * self._PX / 72.0))

    def _wrap(self, d, t, f, maxw):
        out = []; cur = ""
        # For CJK (no spaces) wrap per-character; for Latin wrap per-word.
        units = list(t) if _has_cjk(t) and " " not in t.strip() else t.split(" ")
        joiner = "" if (units and units == list(t)) else " "
        for w in units:
            cand = (cur + joiner + w) if cur else w
            if d.textlength(cand, font=f) <= maxw:
                cur = cand
            else:
                if cur:
                    out.append(cur)
                cur = w
        if cur:
            out.append(cur)
        return out or [""]

    def _render_png(self, spec, path):
        img = Image.new("RGB", (self._I(SW_IN), self._I(SH_IN)), spec["bg"])
        d = ImageDraw.Draw(img)
        for el in spec["el"]:
            kind = el[0]
            if kind == "rect":
                _, x, y, w, h, c = el
                d.rectangle([self._I(x), self._I(y), self._I(x + w), self._I(y + h)], fill=c)
            elif kind == "image":
                _, x, y, w, h, name = el
                ic = self._icons[name].resize((max(1, self._I(w)), max(1, self._I(h))))
                img.paste(ic, (self._I(x), self._I(y)), ic)
            elif kind == "text":
                _, x, y, w, h, t, size, color, bold, italic, align, ls = el
                color = color or self.t.slate
                size *= FONT_SCALE
                cjk = _has_cjk(t)
                f = _pil_font(self._pt(size), bold, italic, cjk)
                maxw = self._I(w); lh = int(self._pt(size) * 1.2 * ls); yy = self._I(y)
                for raw in t.split("\n"):
                    for line in self._wrap(d, raw, f, maxw):
                        tw = d.textlength(line, font=f)
                        xx = self._I(x) + (maxw - tw) // 2 if align == "c" else \
                             self._I(x) + maxw - tw if align == "r" else self._I(x)
                        d.text((xx, yy), line, font=f, fill=color); yy += lh
            elif kind == "bullets":
                _, x, y, w, h, items_, size, color, gap, bold_lead = el
                color = color or self.t.slate
                size *= FONT_SCALE
                maxw = self._I(w); lh = int(self._pt(size) * 1.2); yy = self._I(y)
                for it in items_:
                    cjk = _has_cjk(it)
                    f = _pil_font(self._pt(size), False, False, cjk)
                    fb = _pil_font(self._pt(size), True, False, cjk)
                    bullet = "•  "
                    indent = self._I(x) + d.textlength(bullet, font=f)
                    if bold_lead and ": " in it:
                        lead, rest = it.split(": ", 1)
                        segs = [(bullet + lead + ": ", fb, self.t.navy), (rest, f, color)]
                    else:
                        segs = [(bullet + it, f, color)]
                    # word/character-aware wrapping across styled segments
                    yy = self._draw_wrapped_segments(d, segs, self._I(x), indent, yy, maxw, lh, cjk)
                    yy += int(gap * self._PX / 72.0)
        img.save(path)
        return img

    def _draw_wrapped_segments(self, d, segs, x0, indent, yy, maxw, lh, cjk):
        # Tokenise into (token, font, color); spaces preserved for Latin.
        tokens = []
        for text, f, col in segs:
            if cjk and " " not in text.strip():
                tokens += [(ch, f, col) for ch in text]
            else:
                for i, wd in enumerate(text.split(" ")):
                    tokens.append((wd + (" " if i < len(text.split(" ")) - 1 else " "), f, col))
        line = []; lw = 0; first = True

        def flush(line, yy, first):
            xpos = x0 if first else indent
            for tok, f, col in line:
                d.text((xpos, yy), tok, font=f, fill=col)
                xpos += d.textlength(tok, font=f)

        for tok, f, col in tokens:
            ww = d.textlength(tok, font=f)
            avail = maxw - (0 if first else (indent - x0))
            if lw + ww > avail and line:
                flush(line, yy, first); yy += lh; line = []; lw = 0; first = False
            line.append((tok, f, col)); lw += ww
        if line:
            flush(line, yy, first); yy += lh
        return yy

    def _contact_sheet(self, previews, path):
        imgs = [Image.open(p) for p in previews]
        n = len(imgs)
        cols = 4
        rows = math.ceil(n / cols)
        tw, th = 308, 173
        pad = 12
        sheet = Image.new("RGB", (cols * tw + pad * (cols + 1), rows * th + pad * (rows + 1)), (210, 214, 219))
        for idx, im in enumerate(imgs):
            r, c = divmod(idx, cols)
            sheet.paste(im.resize((tw, th)), (pad + c * (tw + pad), pad + r * (th + pad)))
        sheet.save(path)

    def _build_pdf(self, path, dpi=170):
        old = self._PX
        self._PX = dpi
        _fontcache.clear()
        try:
            import tempfile
            pages = []
            for i, s in enumerate(self.slides):
                pp = os.path.join(tempfile.gettempdir(), f"_pgod_{i + 1}.png")
                self._render_png(s, pp)
                pages.append(Image.open(pp).convert("RGB"))
            pages[0].save(path, "PDF", save_all=True, append_images=pages[1:], resolution=float(dpi))
        finally:
            self._PX = old
            _fontcache.clear()

    # ---- icons (line art, transparent PNG) ----
    def _gen_icons(self):
        col = self.t.accent + (255,)
        SZ = 1000
        sw = 46

        def new():
            im = Image.new("RGBA", (SZ, SZ), (0, 0, 0, 0))
            return im, ImageDraw.Draw(im)

        # globe
        im, d = new(); p = 140; R = (SZ - 2 * p) / 2; cx = cy = SZ / 2
        d.ellipse([p, p, SZ - p, SZ - p], outline=col, width=sw)
        d.line([cx, p, cx, SZ - p], fill=col, width=int(sw * 0.7))
        for rw in (R * 0.42, R * 0.80):
            d.ellipse([cx - rw, p, cx + rw, SZ - p], outline=col, width=int(sw * 0.6))
        for dy in (-R * 0.55, 0, R * 0.55):
            half = math.sqrt(max(R * R - dy * dy, 0))
            d.line([cx - half, cy + dy, cx + half, cy + dy], fill=col, width=int(sw * 0.55))
        self._icons["globe"] = im

        # scales
        im, d = new(); cx = SZ / 2
        d.line([cx, 195, cx, 560], fill=col, width=sw)
        d.line([235, 250, 765, 250], fill=col, width=sw)
        d.ellipse([cx - 32, 150, cx + 32, 214], fill=col)
        d.line([cx, 560, 350, 650], fill=col, width=sw); d.line([cx, 560, 650, 650], fill=col, width=sw)
        d.line([320, 668, 680, 668], fill=col, width=sw)
        for hx in (235, 765):
            d.line([hx, 250, hx - 98, 432], fill=col, width=int(sw * 0.6))
            d.line([hx, 250, hx + 98, 432], fill=col, width=int(sw * 0.6))
            d.arc([hx - 122, 360, hx + 122, 500], 8, 172, fill=col, width=sw)
        self._icons["scales"] = im

        # coins
        im, d = new(); cx = SZ / 2; ew = 460; eh = 156; levels = [610, 475, 340]
        for yc in levels:
            d.ellipse([cx - ew / 2, yc - eh / 2, cx + ew / 2, yc + eh / 2], outline=col, width=sw)
        d.line([cx - ew / 2, levels[-1], cx - ew / 2, levels[0]], fill=col, width=sw)
        d.line([cx + ew / 2, levels[-1], cx + ew / 2, levels[0]], fill=col, width=sw)
        d.ellipse([cx - ew / 2, levels[-1] - eh / 2, cx + ew / 2, levels[-1] + eh / 2], outline=col, width=sw)
        self._icons["coins"] = im

        # doc
        im, d = new(); x0, y0, x1, y1 = 275, 150, 690, 800; fold = 125
        for seg in ([x0, y0, x1 - fold, y0], [x1 - fold, y0, x1, y0 + fold],
                    [x1, y0 + fold, x1, y1], [x1, y1, x0, y1], [x0, y1, x0, y0]):
            d.line(seg, fill=col, width=sw)
        d.line([x1 - fold, y0, x1 - fold, y0 + fold], fill=col, width=int(sw * 0.7))
        d.line([x1 - fold, y0 + fold, x1, y0 + fold], fill=col, width=int(sw * 0.7))
        for ty in (335, 435, 535, 635):
            d.line([x0 + 72, ty, x1 - 72, ty], fill=col, width=int(sw * 0.55))
        self._icons["doc"] = im

        # lightbulb (idea)
        im, d = new(); cx = SZ / 2
        d.ellipse([cx - 230, 150, cx + 230, 610], outline=col, width=sw)
        d.line([cx - 110, 600, cx - 110, 720], fill=col, width=sw)
        d.line([cx + 110, 600, cx + 110, 720], fill=col, width=sw)
        for yy in (720, 780, 840):
            d.line([cx - 120, yy, cx + 120, yy], fill=col, width=int(sw * 0.7))
        d.line([cx, 880, cx, 920], fill=col, width=int(sw * 0.7))
        self._icons["idea"] = im

        # chart (data)
        im, d = new(); base = 800; left = 200
        d.line([left, 160, left, base], fill=col, width=sw)
        d.line([left, base, 840, base], fill=col, width=sw)
        for i, bh in enumerate((220, 360, 300, 470)):
            bx = left + 70 + i * 150
            d.rectangle([bx, base - bh, bx + 95, base], outline=col, width=int(sw * 0.75))
        self._icons["chart"] = im
